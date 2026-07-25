"""
Report export: PDF, Excel, and PowerPoint versions of a dataset's dashboard
(KPIs, trend, breakdowns, insights, recommendations). Generated files are
saved under the top-level reports/ folder and served back for download.
"""
from __future__ import annotations

import os
from datetime import datetime
from typing import Any, Dict

import pandas as pd

from app.ai.recommendations import generate_recommendations
from app.services.analytics import (
    breakdown_by,
    business_insights,
    compute_kpis,
    guess_column_mapping,
    revenue_trend,
)

REPORTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "reports")

# Brand palette (kept in sync with the frontend's Tailwind config)
COLOR_INK = "#123B44"
COLOR_TEAL = "#1F6F7A"
COLOR_GOLD = "#D9A441"


def _report_filename(dataset_name: str, dataset_id: str, ext: str) -> str:
    safe_name = "".join(c if c.isalnum() or c in "-_ " else "_" for c in dataset_name).strip().replace(" ", "_")
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    os.makedirs(REPORTS_DIR, exist_ok=True)
    return os.path.join(REPORTS_DIR, f"{safe_name}_{dataset_id[:8]}_{timestamp}.{ext}")


def _gather_report_data(df: pd.DataFrame) -> Dict[str, Any]:
    mapping = guess_column_mapping(list(df.columns))
    return {
        "mapping": mapping,
        "kpis": compute_kpis(df, mapping),
        "trend": revenue_trend(df, mapping, freq="M"),
        "top_products": breakdown_by(df, mapping, "product", "revenue", top_n=10),
        "top_regions": breakdown_by(df, mapping, "region", "revenue", top_n=10),
        "insights": business_insights(df, mapping),
        "recommendations": generate_recommendations(df),
    }


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def generate_pdf_report(df: pd.DataFrame, dataset_name: str, dataset_id: str) -> str:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak

    data = _gather_report_data(df)
    path = _report_filename(dataset_name, dataset_id, "pdf")

    doc = SimpleDocTemplate(path, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleBrand", parent=styles["Title"], textColor=colors.HexColor(COLOR_INK))
    heading_style = ParagraphStyle("HeadingBrand", parent=styles["Heading2"], textColor=colors.HexColor(COLOR_TEAL))
    body_style = styles["BodyText"]

    story = []
    story.append(Paragraph("InsightFlow AI", title_style))
    story.append(Paragraph(f"Business Report — {dataset_name}", styles["Heading3"]))
    story.append(Paragraph(datetime.utcnow().strftime("Generated %B %d, %Y at %H:%M UTC"), body_style))
    story.append(Spacer(1, 0.3 * inch))

    # KPIs
    story.append(Paragraph("Key Performance Indicators", heading_style))
    kpi_rows = [["Metric", "Value"]] + [
        [k.replace("_", " ").title(), f"{v:,.2f}" if isinstance(v, float) else f"{v:,}"]
        for k, v in data["kpis"].items()
    ]
    kpi_table = Table(kpi_rows, colWidths=[3 * inch, 2.5 * inch])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(COLOR_INK)),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CFE6E6")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#EAF3F3")]),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(kpi_table)
    story.append(Spacer(1, 0.3 * inch))

    # Top products / regions
    for title, rows, label_col in [
        ("Top Products by Revenue", data["top_products"], "Product"),
        ("Top Regions by Revenue", data["top_regions"], "Region"),
    ]:
        if not rows:
            continue
        story.append(Paragraph(title, heading_style))
        table_rows = [[label_col, "Revenue"]] + [[r["label"], f"${r['value']:,.2f}"] for r in rows]
        t = Table(table_rows, colWidths=[3 * inch, 2.5 * inch])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(COLOR_TEAL)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CFE6E6")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#EAF3F3")]),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.25 * inch))

    story.append(PageBreak())

    # Insights
    story.append(Paragraph("Business Insights", heading_style))
    if data["insights"]:
        for key, value in data["insights"].items():
            label = key.replace("_", " ").title()
            display_val = ", ".join(value) if isinstance(value, list) else str(value)
            story.append(Paragraph(f"<b>{label}:</b> {display_val}", body_style))
    else:
        story.append(Paragraph("Not enough recognizable columns to generate insights.", body_style))
    story.append(Spacer(1, 0.3 * inch))

    # Recommendations
    story.append(Paragraph("AI Recommendations", heading_style))
    for rec in data["recommendations"]:
        story.append(Paragraph(f"• {rec}", body_style))

    doc.build(story)
    return path


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def generate_excel_report(df: pd.DataFrame, dataset_name: str, dataset_id: str) -> str:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils.dataframe import dataframe_to_rows

    data = _gather_report_data(df)
    path = _report_filename(dataset_name, dataset_id, "xlsx")

    wb = Workbook()
    header_fill = PatternFill(start_color="123B44", end_color="123B44", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)

    def style_header(ws, row=1):
        for cell in ws[row]:
            cell.fill = header_fill
            cell.font = header_font

    # KPI Summary
    ws = wb.active
    ws.title = "KPI Summary"
    ws.append(["Metric", "Value"])
    for k, v in data["kpis"].items():
        ws.append([k.replace("_", " ").title(), v])
    style_header(ws)
    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 18

    # Trend
    ws2 = wb.create_sheet("Revenue Trend")
    ws2.append(["Period", "Revenue"])
    for point in data["trend"]:
        ws2.append([point["period"], point["revenue"]])
    style_header(ws2)

    # Top products / regions
    ws3 = wb.create_sheet("Top Products")
    ws3.append(["Product", "Revenue"])
    for row in data["top_products"]:
        ws3.append([row["label"], row["value"]])
    style_header(ws3)

    ws4 = wb.create_sheet("Top Regions")
    ws4.append(["Region", "Revenue"])
    for row in data["top_regions"]:
        ws4.append([row["label"], row["value"]])
    style_header(ws4)

    # Raw data (capped to keep file size sane)
    ws5 = wb.create_sheet("Raw Data")
    for r in dataframe_to_rows(df.head(5000), index=False, header=True):
        ws5.append(r)
    style_header(ws5)

    wb.save(path)
    return path


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def generate_pptx_report(df: pd.DataFrame, dataset_name: str, dataset_id: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    data = _gather_report_data(df)
    path = _report_filename(dataset_name, dataset_id, "pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    ink = RGBColor(0x12, 0x3B, 0x44)
    gold = RGBColor(0xD9, 0xA4, 0x41)
    blank_layout = prs.slide_layouts[6]

    def add_title(slide, text, subtitle=None):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ink
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x4B, 0x5A, 0x63)

    # Slide 1: Title
    slide = prs.slides.add_slide(blank_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = box.text_frame
    tf.text = "InsightFlow AI"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ink
    p = tf.add_paragraph()
    p.text = f"Business Report — {dataset_name}"
    p.font.size = Pt(22)
    p.font.color.rgb = gold
    p2 = tf.add_paragraph()
    p2.text = datetime.utcnow().strftime("Generated %B %d, %Y")
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x4B, 0x5A, 0x63)

    # Slide 2: KPIs
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "Key Performance Indicators")
    top, left, col_width = Inches(1.6), Inches(0.6), Inches(3.9)
    for i, (k, v) in enumerate(data["kpis"].items()):
        col = i % 3
        row = i // 3
        box = slide.shapes.add_textbox(left + col * col_width, top + row * Inches(1.6), col_width - Inches(0.2), Inches(1.4))
        tf = box.text_frame
        tf.text = f"{v:,.2f}" if isinstance(v, float) else f"{v:,}"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ink
        p = tf.add_paragraph()
        p.text = k.replace("_", " ").title()
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x4B, 0x5A, 0x63)

    # Slide 3: Revenue trend chart (native pptx chart)
    if data["trend"]:
        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, "Revenue Trend")
        chart_data = CategoryChartData()
        chart_data.categories = [p["period"] for p in data["trend"]]
        chart_data.add_series("Revenue", [p["revenue"] for p in data["trend"]])
        slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), chart_data)

    # Slide 4: Top products chart
    if data["top_products"]:
        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, "Top Products by Revenue")
        chart_data = CategoryChartData()
        chart_data.categories = [r["label"] for r in data["top_products"]]
        chart_data.add_series("Revenue", [r["value"] for r in data["top_products"]])
        slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), chart_data)

    # Slide 5: Recommendations
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "AI Recommendations")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, rec in enumerate(data["recommendations"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"→ {rec}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x2A, 0x36, 0x3C)
        p.space_after = Pt(12)

    prs.save(path)
    return path
